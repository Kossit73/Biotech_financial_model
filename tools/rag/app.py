from __future__ import annotations

import hashlib
import json
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, List, Optional

import faiss
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from fastapi import FastAPI, File, HTTPException, UploadFile
from matplotlib.ticker import FuncFormatter
from pydantic import BaseModel, Field
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from sentence_transformers import CrossEncoder, SentenceTransformer


PROJECT_ROOT = Path(os.environ.get("RAG_PROJECTS_DIR", "projects"))
EMBED_MODEL_NAME = os.environ.get("RAG_EMBED_MODEL", "all-MiniLM-L6-v2")
RERANK_MODEL_NAME = os.environ.get("RAG_RERANK_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")
MAX_CHUNK_CHARS = int(os.environ.get("RAG_CHUNK_CHARS", "1200"))
CHUNK_OVERLAP = int(os.environ.get("RAG_CHUNK_OVERLAP", "200"))


class FinancialSnapshot(BaseModel):
    as_of: Optional[str] = None
    workbook_path: Optional[str] = None
    workbook_hash: Optional[str] = None
    currency: Optional[str] = None
    assumptions: dict = Field(default_factory=dict)
    capex_total: Optional[float] = None
    opex_annual: Optional[float] = None
    revenue_annual: Optional[float] = None
    npv: Optional[float] = None
    irr: Optional[float] = None
    payback_years: Optional[float] = None
    dscr_min: Optional[float] = None
    sensitivities: list = Field(default_factory=list)
    scenarios: list = Field(default_factory=list)


class CollectRequest(BaseModel):
    project_id: str
    financial_snapshot: FinancialSnapshot
    cell_map: dict = Field(default_factory=dict)
    workbook_hash: Optional[str] = None


class GenerateRequest(BaseModel):
    project_id: str
    section_outline: Optional[List[str]] = None


class ChunkMetadata(BaseModel):
    project_id: str
    file_path: str
    file_type: str
    page_or_sheet: str
    section: Optional[str]
    char_start: int
    char_end: int
    hash: str


app = FastAPI()


def _project_dir(project_id: str) -> Path:
    base = PROJECT_ROOT / project_id
    base.mkdir(parents=True, exist_ok=True)
    (base / "uploads").mkdir(exist_ok=True)
    (base / "parsed").mkdir(exist_ok=True)
    (base / "index").mkdir(exist_ok=True)
    (base / "financial").mkdir(exist_ok=True)
    (base / "reports").mkdir(exist_ok=True)
    (base / "charts").mkdir(exist_ok=True)
    return base


def _hash_file(path: Path) -> str:
    sha = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            sha.update(chunk)
    return sha.hexdigest()


def _stream_upload(upload: UploadFile, destination: Path) -> None:
    with destination.open("wb") as handle:
        while True:
            chunk = upload.file.read(1024 * 1024)
            if not chunk:
                break
            handle.write(chunk)


def _chunk_text(text: str, max_chars: int = MAX_CHUNK_CHARS, overlap: int = CHUNK_OVERLAP) -> List[str]:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if not cleaned:
        return []
    chunks = []
    start = 0
    while start < len(cleaned):
        end = min(start + max_chars, len(cleaned))
        chunks.append(cleaned[start:end])
        start = max(end - overlap, 0)
    return chunks


def _extract_pdf(path: Path) -> Iterable[tuple[str, str]]:
    reader = PdfReader(str(path))
    for idx, page in enumerate(reader.pages, start=1):
        yield str(idx), page.extract_text() or ""


def _extract_docx(path: Path) -> Iterable[tuple[str, str]]:
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    yield "document", "\n".join(paragraphs)


def _extract_pptx(path: Path) -> Iterable[tuple[str, str]]:
    presentation = Presentation(str(path))
    for idx, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        yield f"slide-{idx}", "\n".join(texts)


def _extract_text(path: Path) -> Iterable[tuple[str, str]]:
    yield "text", path.read_text(encoding="utf-8", errors="ignore")


def _extract_csv(path: Path) -> Iterable[tuple[str, str]]:
    df = pd.read_csv(path)
    yield "csv", df.to_csv(index=False)


def _extract_xlsx(path: Path) -> Iterable[tuple[str, str]]:
    sheets = pd.read_excel(path, sheet_name=None)
    for name, df in sheets.items():
        yield name, df.to_csv(index=False)


def _extract_content(path: Path) -> Iterable[tuple[str, str]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".csv":
        return _extract_csv(path)
    if suffix in {".xlsx", ".xls"}:
        return _extract_xlsx(path)
    return _extract_text(path)


def _save_metadata(index_dir: Path, metadata: List[ChunkMetadata]) -> None:
    meta_path = index_dir / "meta.jsonl"
    with meta_path.open("a", encoding="utf-8") as handle:
        for entry in metadata:
            handle.write(entry.model_dump_json())
            handle.write("\n")


def _load_metadata(index_dir: Path) -> List[ChunkMetadata]:
    meta_path = index_dir / "meta.jsonl"
    if not meta_path.exists():
        return []
    entries = []
    with meta_path.open("r", encoding="utf-8") as handle:
        for line in handle:
            entries.append(ChunkMetadata.model_validate_json(line))
    return entries


def _load_index(index_dir: Path, dim: int) -> faiss.Index:
    index_path = index_dir / "index.faiss"
    if index_path.exists():
        return faiss.read_index(str(index_path))
    return faiss.IndexFlatIP(dim)


def _persist_index(index_dir: Path, index: faiss.Index) -> None:
    faiss.write_index(index, str(index_dir / "index.faiss"))


def _embed_texts(model: SentenceTransformer, texts: List[str]) -> np.ndarray:
    embeddings = model.encode(texts, normalize_embeddings=True)
    return np.asarray(embeddings, dtype="float32")


def _serialize_snapshot(snapshot: FinancialSnapshot, cell_map: dict, workbook_hash: Optional[str]) -> dict:
    payload = snapshot.model_dump()
    payload["as_of"] = snapshot.as_of or datetime.now(timezone.utc).isoformat()
    payload["workbook_hash"] = workbook_hash or snapshot.workbook_hash
    payload["cell_map"] = cell_map
    return payload


def _financial_bullets(snapshot: dict) -> str:
    bullets = []
    for key in [
        "npv",
        "irr",
        "dscr_min",
        "payback_years",
        "capex_total",
        "opex_annual",
        "revenue_annual",
    ]:
        if snapshot.get(key) is not None:
            bullets.append(f"- {key.replace('_', ' ').title()}: {snapshot[key]}")
    if snapshot.get("scenarios"):
        bullets.append("- Scenarios:")
        for scenario in snapshot["scenarios"]:
            bullets.append(f"  - {scenario.get('name')}: NPV {scenario.get('npv')}, IRR {scenario.get('irr')}")
    if snapshot.get("sensitivities"):
        bullets.append("- Sensitivities:")
        for sensitivity in snapshot["sensitivities"]:
            bullets.append(
                f"  - {sensitivity.get('variable')} {sensitivity.get('delta')}: "
                f"NPV {sensitivity.get('npv')}, IRR {sensitivity.get('irr')}"
            )
    return "\n".join(bullets)


def _section_outline() -> List[str]:
    return [
        "Executive Summary",
        "Project Description & Scope",
        "Market & Demand Analysis",
        "Technical & Operations",
        "Legal, Permitting & Environmental",
        "Implementation Plan",
        "Financial Analysis",
        "Risk Assessment & Mitigations",
        "Conclusion & Recommendation",
        "Appendices",
    ]


def _section_queries() -> dict:
    return {
        "Executive Summary": "materiality of results, decision drivers, showstoppers",
        "Market & Demand Analysis": "market size, demand forecast, price assumptions, offtake",
        "Technical & Operations": "process design, throughput, yield, utilities, site layout",
        "Legal, Permitting & Environmental": "permits, EIA/ESIA, land rights, community",
        "Implementation Plan": "schedule, capex phasing, procurement strategy, org",
        "Financial Analysis": "NPV, IRR, DSCR, payback, sensitivities",
        "Risk Assessment & Mitigations": "risk register, mitigation, ESG metrics",
    }


def _ensure_chart_dir(project_dir: Path) -> Path:
    chart_dir = project_dir / "charts"
    chart_dir.mkdir(parents=True, exist_ok=True)
    return chart_dir


def _fmt_currency(curr: str) -> FuncFormatter:
    def _fmt(x, pos):
        try:
            return f"{curr} {x:,.0f}"
        except Exception:
            return f"{x:,.0f}"
    return FuncFormatter(_fmt)


def plot_npv_curve(financial: dict, out_path: Path) -> Optional[Path]:
    xs: List[str] = []
    ys: List[float] = []
    if financial.get("scenarios"):
        xs = [s.get("name", f"S{i + 1}") for i, s in enumerate(financial["scenarios"])]
        ys = [float(s.get("npv", 0)) for s in financial["scenarios"]]
    elif financial.get("sensitivities"):
        sens = [
            s for s in financial["sensitivities"]
            if s.get("variable") == "price" and s.get("npv") is not None
        ]
        sens = sorted(sens, key=lambda s: s.get("delta", 0))
        xs = [f"{int(s.get('delta', 0) * 100)}%" for s in sens]
        ys = [float(s["npv"]) for s in sens]
    else:
        return None

    fig = plt.figure()
    plt.plot(xs, ys, marker="o")
    plt.title("NPV Curve")
    plt.xlabel("Scenario / Delta")
    plt.ylabel("NPV")
    curr = financial.get("currency", "") or ""
    plt.gca().yaxis.set_major_formatter(_fmt_currency(curr))
    plt.grid(True, linestyle="--", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out_path, dpi=200)
    plt.close(fig)
    return out_path


def plot_dscr_trend_from_excel(
    xlsx_path: str,
    sheet: str,
    date_col: str,
    dscr_col: str,
    out_path: Path,
) -> Optional[Path]:
    if not xlsx_path:
        return None
    df = pd.read_excel(xlsx_path, sheet_name=sheet)
    if date_col not in df.columns or dscr_col not in df.columns:
        return None
    df = df[[date_col, dscr_col]].dropna()
    df[date_col] = pd.to_datetime(df[date_col], errors="coerce")
    df = df.dropna(subset=[date_col, dscr_col])

    fig = plt.figure()
    plt.plot(df[date_col], df[dscr_col])
    plt.title("DSCR Trend")
    plt.xlabel("Date")
    plt.ylabel("DSCR")
    plt.axhline(1.0, linestyle="--", linewidth=1)
    plt.grid(True, linestyle="--", alpha=0.4)
    fig.autofmt_xdate()
    fig.tight_layout()
    fig.savefig(out_path, dpi=200)
    plt.close(fig)
    return out_path


def _compose_section(section: str, financial_snapshot: dict, passages: List[dict]) -> str:
    lines = [f"## {section}", "", "[FINANCIAL_SNAPSHOT]", _financial_bullets(financial_snapshot), ""]
    lines.append("[CONTEXT]")
    for passage in passages:
        citation = f"[Source: {passage['file']} {passage['page']}]"
        lines.append(f"- {passage['text']} {citation}")
    lines.append("")
    return "\n".join(lines)


def _retrieve_passages(
    project_dir: Path,
    embed_model: SentenceTransformer,
    query: str,
    top_k: int = 5,
) -> List[dict]:
    index_dir = project_dir / "index"
    metadata = _load_metadata(index_dir)
    if not metadata:
        return []
    dim = embed_model.get_sentence_embedding_dimension()
    index = _load_index(index_dir, dim)
    query_emb = _embed_texts(embed_model, [query])
    scores, idxs = index.search(query_emb, min(top_k, len(metadata)))
    passages = []
    for idx in idxs[0]:
        meta = metadata[idx]
        parsed_text = (project_dir / "parsed" / f"{Path(meta.file_path).stem}.json").read_text()
        parsed_data = json.loads(parsed_text)
        chunk_text = parsed_data[str(meta.page_or_sheet)][meta.char_start:meta.char_end]
        passages.append(
            {
                "file": Path(meta.file_path).name,
                "page": meta.page_or_sheet,
                "text": chunk_text,
            }
        )
    return passages


@app.post("/collect")
def collect_financials(payload: CollectRequest) -> dict:
    project_dir = _project_dir(payload.project_id)
    snapshot = _serialize_snapshot(payload.financial_snapshot, payload.cell_map, payload.workbook_hash)
    snapshot_path = project_dir / "financial" / "snapshot.json"
    snapshot_path.write_text(json.dumps(snapshot, indent=2))
    return {"status": "ok", "snapshot_path": str(snapshot_path)}


@app.post("/ingest")
def ingest_files(project_id: str, files: List[UploadFile] = File(...)) -> dict:
    project_dir = _project_dir(project_id)
    embed_model = SentenceTransformer(EMBED_MODEL_NAME)
    index_dir = project_dir / "index"
    dim = embed_model.get_sentence_embedding_dimension()
    index = _load_index(index_dir, dim)
    all_metadata: List[ChunkMetadata] = []

    for upload in files:
        suffix = Path(upload.filename).suffix.lower()
        dest = project_dir / "uploads" / f"{uuid.uuid4().hex}{suffix}"
        _stream_upload(upload, dest)
        file_hash = _hash_file(dest)
        parsed_texts = {}
        for page, text in _extract_content(dest):
            parsed_texts[str(page)] = text
            chunks = _chunk_text(text)
            if not chunks:
                continue
            embeddings = _embed_texts(embed_model, chunks)
            index.add(embeddings)
            start = 0
            for chunk in chunks:
                end = start + len(chunk)
                all_metadata.append(
                    ChunkMetadata(
                        project_id=project_id,
                        file_path=str(dest),
                        file_type=suffix.lstrip("."),
                        page_or_sheet=str(page),
                        section=None,
                        char_start=start,
                        char_end=end,
                        hash=file_hash,
                    )
                )
                start = max(end - CHUNK_OVERLAP, 0)
        parsed_path = project_dir / "parsed" / f"{dest.stem}.json"
        parsed_path.write_text(json.dumps(parsed_texts, indent=2))

    _persist_index(index_dir, index)
    _save_metadata(index_dir, all_metadata)
    return {"status": "ok", "chunks_added": len(all_metadata)}


@app.post("/generate")
def generate_report(payload: GenerateRequest) -> dict:
    project_dir = _project_dir(payload.project_id)
    snapshot_path = project_dir / "financial" / "snapshot.json"
    if not snapshot_path.exists():
        raise HTTPException(status_code=404, detail="Financial snapshot not found.")
    financial_snapshot = json.loads(snapshot_path.read_text())
    outline = payload.section_outline or _section_outline()
    embed_model = SentenceTransformer(EMBED_MODEL_NAME)
    reranker = CrossEncoder(RERANK_MODEL_NAME)
    chart_dir = _ensure_chart_dir(project_dir)
    npv_png = plot_npv_curve(financial_snapshot, chart_dir / "npv_curve.png")
    dscr_png = plot_dscr_trend_from_excel(
        financial_snapshot.get("workbook_path", ""),
        sheet="Debt",
        date_col="Date",
        dscr_col="DSCR",
        out_path=chart_dir / "dscr_trend.png",
    )

    sections = []
    for section in outline:
        query = _section_queries().get(section, section)
        passages = _retrieve_passages(project_dir, embed_model, query, top_k=8)
        if passages:
            pairs = [(query, p["text"]) for p in passages]
            scores = reranker.predict(pairs)
            passages = [p for _, p in sorted(zip(scores, passages), key=lambda pair: pair[0], reverse=True)]
        sections.append(_compose_section(section, financial_snapshot, passages[:5]))

    report = "\n\n".join(sections)
    chart_lines = []
    if npv_png and npv_png.exists():
        chart_lines.append("### NPV Curve")
        chart_lines.append(f"![NPV Curve]({npv_png.as_posix()})")
    if dscr_png and dscr_png.exists():
        chart_lines.append("### DSCR Trend")
        chart_lines.append(f"![DSCR Trend]({dscr_png.as_posix()})")
    if chart_lines:
        report = f"{report}\n\n## Charts\n\n" + "\n".join(chart_lines) + "\n"
    report_path = project_dir / "reports" / "report.md"
    report_path.write_text(report)
    return {"status": "ok", "report_path": str(report_path)}
